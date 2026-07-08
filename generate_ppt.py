"""Generate OOA Microsite Workflow PPT using python-pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

NAVY  = RGBColor(0x04, 0x04, 0x3a)
GOLD  = RGBColor(0xC5, 0xA5, 0x72)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT= RGBColor(0x0A, 0x2A, 0x44)
GREY  = RGBColor(0x64, 0x74, 0x8B)
GREEN = RGBColor(0x06, 0x5F, 0x46)

W, H = Inches(13.33), Inches(7.5)  # 16:9
prs = Presentation()
prs.slide_width = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # completely blank

def add_rect(slide, l, t, w, h, fill=None, line=None, alpha=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background() if not line else None
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if not line:
        s.line.fill.background()
    return s

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return tb

def slide_header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
    add_rect(slide, 0, 0, 13.33, 1.6, ACCENT)
    add_rect(slide, 0, 1.55, 13.33, 0.08, GOLD)
    add_text(slide, title, 0.4, 0.2, 12.5, 1.1, size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 1.1, 12.5, 0.5, size=14, color=GOLD, align=PP_ALIGN.LEFT)

def add_bullet_box(slide, items, l, t, w, h, col=WHITE):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"  ▸  {item}"
        run.font.size = Pt(13)
        run.font.color.rgb = col
        run.font.name = 'Calibri'

def feature_box(slide, l, t, w, h, icon, title, desc, bg=ACCENT):
    add_rect(slide, l, t, w, h, bg)
    add_text(slide, icon, l+0.1, t+0.1, 0.5, 0.4, size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, title, l+0.55, t+0.08, w-0.65, 0.35, size=12, bold=True, color=WHITE)
    add_text(slide, desc, l+0.55, t+0.42, w-0.65, h-0.5, size=10, color=LIGHT)

def step_row(slide, num, title, desc, t):
    add_rect(slide, 0.4, t, 0.5, 0.5, GOLD)
    add_text(slide, str(num), 0.4, t, 0.5, 0.5, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, title, 1.05, t, 3.5, 0.25, size=12, bold=True, color=GOLD)
    add_text(slide, desc, 1.05, t+0.25, 11.8, 0.28, size=10, color=LIGHT)
    add_rect(slide, 0.4, t+0.52, 12.5, 0.02, GREY)

# ════ SLIDE 1: COVER ════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, NAVY)
add_rect(sl, 0, 0, 13.33, 0.15, GOLD)
add_rect(sl, 0, 7.35, 13.33, 0.15, GOLD)
# Big title
add_text(sl, "OOA JAIN University", 1, 1.2, 11.3, 1.2, size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Digital Platform — Workflow & Feature Guide", 1, 2.4, 11.3, 0.8, size=22, bold=False, color=GOLD, align=PP_ALIGN.CENTER)
add_rect(sl, 3, 3.3, 7.33, 0.05, GOLD)
add_text(sl, "Office of Academics  |  Jain (Deemed-to-be University)", 1, 3.5, 11.3, 0.6, size=16, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(sl, "Version 1.0  |  2026", 1, 4.1, 11.3, 0.5, size=14, color=GREY, align=PP_ALIGN.CENTER)
add_text(sl, "github.com/ooa-jain/ooa-microsite", 1, 6.5, 11.3, 0.5, size=13, color=GOLD, align=PP_ALIGN.CENTER)

# ════ SLIDE 2: AGENDA ════
sl = prs.slides.add_slide(blank)
slide_header(sl, "Agenda", "What this presentation covers")
items = ["01  Platform Overview & Architecture",
         "02  Landing Page — Login & Registration",
         "03  Home Dashboard — All Panels",
         "04  Academic Calendar (2800+ Events)",
         "05  UGC & University Notices",
         "06  Event Reminders & Email Alerts",
         "07  User Profile Management",
         "08  Document Repositories",
         "09  15-Minute Course & Guided Tour",
         "10  Admin Dashboard & Controls",
         "11  Complete Workflow Summary"]
add_bullet_box(sl, items[:6], 0.6, 1.8, 6, 5)
add_bullet_box(sl, items[6:], 6.8, 1.8, 6, 5)

# ════ SLIDE 3: ARCHITECTURE ════
sl = prs.slides.add_slide(blank)
slide_header(sl, "Platform Architecture", "Technology stack powering the OOA Microsite")
boxes = [
    (0.3,  2.0, 2.8, 1.6, "🐍", "Backend", "Python 3.13 + Flask\nJinja2 Templates"),
    (3.3,  2.0, 2.8, 1.6, "🍃", "Database", "MongoDB Atlas\n27 Collections"),
    (6.3,  2.0, 2.8, 1.6, "🔐", "Auth", "Flask Sessions\n+ Firebase OAuth"),
    (9.3,  2.0, 3.5, 1.6, "📧", "Email", "SMTP HTML Emails\nAuto Reminders"),
    (0.3,  3.9, 2.8, 1.6, "🎨", "Frontend", "HTML + CSS + JS\nNo framework"),
    (3.3,  3.9, 2.8, 1.6, "🚀", "Deploy", "Render.com\nPort 5001"),
    (6.3,  3.9, 2.8, 1.6, "👥", "Roles", "Admin / User\nLeader / Core"),
    (9.3,  3.9, 3.5, 1.6, "📁", "Storage", "Local uploads/\nMongoDB GridFS"),
]
for l,t,w,h,ic,tl,ds in boxes:
    feature_box(sl, l, t, w, h, ic, tl, ds)

# ════ SLIDE 4: LANDING PAGE ════
sl = prs.slides.add_slide(blank)
slide_header(sl, "Landing Page  ( / )", "Public entry point with login, Google SSO, and registration")
add_text(sl, "What users see first:", 0.5, 1.8, 12, 0.4, size=14, bold=True, color=GOLD)
steps = [
    (1, "Video Hero", "Full-screen dark video with 'Office of Academics' title, 28+ departments stat"),
    (2, "Login Card", "Email + Password form, Google Sign-In (Jain domain only), Forgot Password link"),
    (3, "Register", "New users register → Admin approves → Access granted to dashboard"),
    (4, "Redirect", "On success: redirected to /home with full dashboard access"),
]
for n,ti,de in steps:
    step_row(sl, n, ti, de, 2.25 + (n-1)*1.1)

# ════ SLIDE 5: HOME DASHBOARD ════
sl = prs.slides.add_slide(blank)
slide_header(sl, "Home Dashboard  ( /home )", "Central hub — loaded asynchronously from /api/home_data")
panels = [
    (0.3,  1.8, 4.0, 1.4, "📢", "Events Ticker", "Today's events scroll slowly across the top"),
    (4.5,  1.8, 4.0, 1.4, "📋", "UGC Notices", "Latest UGC circulars & guidelines panel"),
    (8.7,  1.8, 4.3, 1.4, "🏛️", "OOA Notices", "Office of Academics announcements panel"),
    (0.3,  3.4, 4.0, 1.4, "📅", "Calendar Slide", "Interactive academic events with dept filter"),
    (4.5,  3.4, 4.0, 1.4, "🖼️", "Campus Gallery", "Photo moments from campus events"),
    (8.7,  3.4, 4.3, 1.4, "📖", "Mini Course", "15-min guided course on JAIN & platform"),
    (0.3,  5.0, 4.0, 1.4, "🎯", "Events Grid", "Today + Upcoming with Set Reminder buttons"),
    (4.5,  5.0, 4.0, 1.4, "🗺️", "Guided Tour", "Driver.js overlay tour of all sections"),
    (8.7,  5.0, 4.3, 1.4, "👤", "User Menu", "Profile, settings, repo links, logout"),
]
for l,t,w,h,ic,tl,ds in panels:
    feature_box(sl, l, t, w, h, ic, tl, ds, ACCENT)

# ════ SLIDE 6: ACADEMIC CALENDAR ════
sl = prs.slides.add_slide(blank)
slide_header(sl, "Academic Calendar  ( /calendar )", "2800+ events for 2026-27 — fully interactive")
add_bullet_box(sl, [
    "Month View  —  Click any day to see events starting or ongoing that day",
    "Timeline View  —  Browse events chronologically month by month",
    "Today View  —  See only events active on today's date",
    "Upcoming View  —  Next 60 days across all departments",
    "Sidebar Navigation  —  Faculty → Department → Sub-department tree with event counts",
    "Campus Filter  —  Toggle between Bengaluru and Kochi campus events",
    "Event Drawer  —  Click any event for full details: venue, coordinator, dates, Add to Google Calendar",
    "Search  —  Real-time keyword search across all 2800+ events",
    "Deep Link  —  /calendar?eventName=X opens specific event drawer directly (from ticker/notices)",
], 0.5, 1.8, 12.4, 5.5, col=LIGHT)

# ════ SLIDE 7: NOTICES ════
sl = prs.slides.add_slide(blank)
slide_header(sl, "UGC & University Notices", "4 UGC + 4 OOA notices — real external PDF links")
add_text(sl, "UGC Notices  (category = 'UGC')", 0.5, 1.8, 6, 0.4, size=13, bold=True, color=GOLD)
add_bullet_box(sl, [
    "UGC Guidelines: DIAC, DPAC and BOS",
    "National Credit Framework (NCrF) – 2023",
    "Reservation in Temporary Appointments",
    "UGC Handbook: Adoption of NEP 2020",
], 0.5, 2.3, 6, 2.5, col=LIGHT)

add_text(sl, "OOA / University Notices  (category = 'University Notice')", 6.8, 1.8, 6.2, 0.4, size=13, bold=True, color=GOLD)
add_bullet_box(sl, [
    "JAIN Academic Calendar 2026-27",
    "Semester-End Examination Schedule",
    "Academic Bank of Credits (ABC) Registration",
    "CBCS Implementation – Updated Structure",
], 6.8, 2.3, 6.2, 2.5, col=LIGHT)

add_rect(sl, 0.5, 4.9, 12.3, 0.05, GOLD)
add_text(sl, "How it works: Click any notice card → Modal opens with description + external PDF link → View or download", 0.5, 5.0, 12.3, 0.6, size=12, color=LIGHT)
add_text(sl, "Admin adds via /edit_ugc — select category 'UGC' or 'University Notice' when uploading", 0.5, 5.6, 12.3, 0.5, size=11, color=GREY)

# ════ SLIDE 8: REMINDERS ════
sl = prs.slides.add_slide(blank)
slide_header(sl, "Event Reminders & Email Alerts", "Set personalised reminders for any academic event")
steps = [
    (1, "Click Event Card", "Click any event in today/upcoming section on Home dashboard"),
    (2, "Reminder Modal Opens", "See event name, date, time, venue in a beautiful modal"),
    (3, "Set Date & Time", "Pick reminder date, time, and optional personal note"),
    (4, "Confirmation Email", "Instant HTML email sent confirming your reminder is set"),
    (5, "Automated Reminder", "At scheduled time, reminder email sent to your Jain email"),
    (6, "Manage in Profile", "View all reminders in My Profile → Cancel any time with one click"),
]
for n,ti,de in steps:
    step_row(sl, n, ti, de, 1.8 + (n-1)*0.93)

# ════ SLIDE 9: PROFILE ════
sl = prs.slides.add_slide(blank)
slide_header(sl, "User Profile  ( /profile )", "Personal info, security, reminders and documents")
boxes = [
    (0.3,  1.8, 3.9, 2.5, "🖼️", "Profile Card", "Name, email, role, dept, location\nApproval status badge"),
    (4.4,  1.8, 3.9, 2.5, "📷", "Profile Photo", "Click avatar → Upload photo\nJPG/PNG/WEBP supported"),
    (8.5,  1.8, 4.5, 2.5, "🔒", "Change Password", "Current → New → Confirm\nServer-side validation"),
    (0.3,  4.5, 3.9, 2.5, "🔔", "My Reminders", "Table: event, date, reminder time\nCancel button per row"),
    (4.4,  4.5, 3.9, 2.5, "📄", "My Documents", "All uploaded docs across\noffice, dept, public repos"),
    (8.5,  4.5, 4.5, 2.5, "✅", "Doc Status", "Approved / Pending / Uploaded\nColour-coded status pills"),
]
for l,t,w,h,ic,tl,ds in boxes:
    feature_box(sl, l, t, w, h, ic, tl, ds)

# ════ SLIDE 10: REPOSITORIES ════
sl = prs.slides.add_slide(blank)
slide_header(sl, "Document Repositories", "Secure document management across three levels")
add_bullet_box(sl, [
    "📁  Department Repo  ( /dept_repo )  —  Faculty upload docs for core team review. Dept-specific access.",
    "🏛️  University Repo  ( /university_repo )  —  Docs shared by core/leader to All Departments.",
    "⚙️  Core Repo  ( /core_repo )  —  Leaders/core upload and target docs to specific departments.",
    "👁️  Document Viewer  ( /view_doc/<id> )  —  Embedded PDF viewer for approved documents.",
    "📏  File Validation  —  Allowed types: PDF, DOCX, images, Excel | Max 50MB per file.",
    "📅  Expiry Dates  —  Documents hidden after end_date passes — automatic cleanup.",
    "🔐  Security Level  —  view_only vs download permissions configurable per document.",
], 0.5, 1.8, 12.5, 5.5, col=LIGHT)

# ════ SLIDE 11: COURSE & TOUR ════
sl = prs.slides.add_slide(blank)
slide_header(sl, "15-Minute Course & Guided Tour", "Built into the Home dashboard for all users")
add_text(sl, "Mini Course Slides:", 0.5, 1.8, 6, 0.4, size=13, bold=True, color=GOLD)
add_bullet_box(sl, [
    "Slide 1 — Welcome to OOA Platform",
    "Slide 2 — About JAIN University (founded 1990)",
    "Slide 3 — NEP 2020 Implementation at JAIN",
    "Slide 4 — Platform Features Walkthrough",
    "Slide 5 — How to Set Reminders",
    "Slide 6 — Quiz (3 questions with instant feedback)",
], 0.5, 2.3, 6, 3.5, col=LIGHT)

add_text(sl, "Guided Product Tour (Driver.js):", 7.0, 1.8, 6, 0.4, size=13, bold=True, color=GOLD)
add_bullet_box(sl, [
    "Step 1 — Today's Events Ticker",
    "Step 2 — UGC Notices Panel",
    "Step 3 — OOA Announcements Panel",
    "Step 4 — Academic Calendar Section",
    "Step 5 — Events Grid with Reminders",
    "Step 6 — Campus Gallery",
    "Step 7 — User Profile & Dashboard",
], 7.0, 2.3, 6, 3.5, col=LIGHT)
add_text(sl, "Click 'Tour' in navbar to launch  •  Press Escape or Done to exit at any time", 0.5, 6.0, 12.5, 0.5, size=11, color=GREY, align=PP_ALIGN.CENTER)

# ════ SLIDE 12: ADMIN ════
sl = prs.slides.add_slide(blank)
slide_header(sl, "Admin Dashboard  ( /admin_dashboard )", "Full control over users, events, and content")
boxes = [
    (0.3,  1.8, 3.9, 1.6, "👥", "User Management", "Approve, role-change, impersonate, reset passwords"),
    (4.4,  1.8, 3.9, 1.6, "📅", "Events Control", "Add/edit/delete academic calendar events"),
    (8.5,  1.8, 4.5, 1.6, "📢", "UGC Notices", "Upload UGC/University notices with files"),
    (0.3,  3.6, 3.9, 1.6, "📰", "Newsletter", "Create and send HTML newsletters"),
    (4.4,  3.6, 3.9, 1.6, "🖼️", "Campus Gallery", "Add/remove campus photo moments"),
    (8.5,  3.6, 4.5, 1.6, "📊", "Activity Logs", "Per-user activity tracking and reports"),
    (0.3,  5.4, 3.9, 1.6, "🏛️", "Campus Setup", "Add/edit campus locations & metadata"),
    (4.4,  5.4, 3.9, 1.6, "📋", "Monthly Engage", "Upload monthly engagement records"),
    (8.5,  5.4, 4.5, 1.6, "⚙️", "ERP Portal", "Semester Readiness & Closure modules"),
]
for l,t,w,h,ic,tl,ds in boxes:
    feature_box(sl, l, t, w, h, ic, tl, ds)

# ════ SLIDE 13: WORKFLOW SUMMARY ════
sl = prs.slides.add_slide(blank)
slide_header(sl, "Complete User Workflow", "End-to-end journey on the OOA platform")
steps = [
    (1, "Visit Landing Page", "Video hero + Login card at / — first impression"),
    (2, "Login / Register", "Credentials or Google SSO → Admin approves new users"),
    (3, "Home Dashboard", "Ticker, Notices, Calendar, Gallery all load from /api/home_data"),
    (4, "Explore Calendar", "2800+ events, filter by dept/campus, view in Month/Timeline/Today view"),
    (5, "Set Reminders", "Click event → Modal → Pick time → Confirmation email sent"),
    (6, "Read Notices", "UGC & OOA panels → Click → Modal with description + PDF link"),
    (7, "Take Course", "15-min slides + quiz about JAIN University and platform"),
    (8, "Manage Profile", "Upload photo, change password, view all reminders & documents"),
]
for n,ti,de in steps:
    step_row(sl, n, ti, de, 1.8 + (n-1)*0.7)

# ════ SLIDE 14: THANK YOU ════
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, NAVY)
add_rect(sl, 0, 0, 13.33, 0.12, GOLD)
add_rect(sl, 0, 7.38, 13.33, 0.12, GOLD)
add_text(sl, "Thank You", 1, 2.0, 11.3, 1.2, size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "OOA Jain University — Office of Academics", 1, 3.3, 11.3, 0.6, size=18, color=GOLD, align=PP_ALIGN.CENTER)
add_rect(sl, 3, 4.0, 7.33, 0.05, GOLD)
add_text(sl, "🌐  github.com/ooa-jain/ooa-microsite", 1, 4.15, 11.3, 0.5, size=14, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(sl, "✉️  ooa.connect@jainuniversity.ac.in", 1, 4.65, 11.3, 0.5, size=14, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(sl, "© 2026 JAIN (Deemed-to-be University). All rights reserved.", 1, 6.8, 11.3, 0.5, size=11, color=GREY, align=PP_ALIGN.CENTER)

prs.save("OOA_Jain_Workflow.pptx")
print("PPT created: OOA_Jain_Workflow.pptx")
